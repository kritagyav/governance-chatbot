"""
Export chat conversation to Word, Excel, PowerPoint, and PDF.
"""

import io
from datetime import datetime


def _qa_pairs(messages: list[dict]) -> list[tuple[dict, dict]]:
    """Extract (user, assistant) pairs from message history."""
    pairs = []
    for i in range(0, len(messages) - 1, 2):
        u, a = messages[i], messages[i + 1]
        if u.get("role") == "user" and a.get("role") == "assistant":
            pairs.append((u, a))
    return pairs


def to_docx(messages: list[dict]) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.2)
        section.right_margin = Inches(1.2)

    # Title
    title = doc.add_heading("Governance Assistant", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.runs[0]
    run.font.color.rgb = RGBColor(0x1B, 0x3A, 0x6B)

    sub = doc.add_paragraph(f"Conversation Export  ·  {datetime.now().strftime('%d %B %Y, %H:%M')}")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].font.size = Pt(10)
    sub.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    doc.add_paragraph()

    pairs = _qa_pairs(messages)
    for i, (q, a) in enumerate(pairs, 1):
        # Question
        q_para = doc.add_paragraph()
        q_label = q_para.add_run(f"Q{i}.  ")
        q_label.bold = True
        q_label.font.color.rgb = RGBColor(0x1B, 0x3A, 0x6B)
        q_text = q_para.add_run(q["content"])
        q_text.bold = True
        q_para.paragraph_format.space_before = Pt(12)

        # Answer
        a_para = doc.add_paragraph(a["content"])
        a_para.paragraph_format.left_indent = Inches(0.3)
        a_para.paragraph_format.space_before = Pt(4)

        # Sources
        sources = a.get("sources", [])
        if sources:
            s_para = doc.add_paragraph("Sources: " + ", ".join(sources))
            s_para.paragraph_format.left_indent = Inches(0.3)
            s_run = s_para.runs[0]
            s_run.italic = True
            s_run.font.size = Pt(9)
            s_run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

        # Divider
        doc.add_paragraph().paragraph_format.space_after = Pt(2)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def to_xlsx(messages: list[dict]) -> bytes:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Conversation"

    navy = "1B3A6B"
    light_blue = "EFF6FF"
    border_color = "CBD5E1"

    thin = Side(style="thin", color=border_color)
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    # Header row
    headers = ["#", "Question", "Answer", "Sources", "Timestamp"]
    col_widths = [5, 45, 70, 35, 20]
    for col, (h, w) in enumerate(zip(headers, col_widths), 1):
        cell = ws.cell(row=1, column=col, value=h)
        cell.font = Font(bold=True, color="FFFFFF", size=11)
        cell.fill = PatternFill("solid", fgColor=navy)
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border
        ws.column_dimensions[cell.column_letter].width = w
    ws.row_dimensions[1].height = 28

    pairs = _qa_pairs(messages)
    ts = datetime.now().strftime("%d %b %Y %H:%M")
    for i, (q, a) in enumerate(pairs, 1):
        row = i + 1
        fill = PatternFill("solid", fgColor=light_blue) if i % 2 == 0 else PatternFill("solid", fgColor="FFFFFF")

        values = [
            i,
            q["content"],
            a["content"],
            ", ".join(a.get("sources", [])),
            ts,
        ]
        for col, val in enumerate(values, 1):
            cell = ws.cell(row=row, column=col, value=val)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = border
            cell.fill = fill
            if col == 1:
                cell.alignment = Alignment(horizontal="center", vertical="top")
            if col == 4:
                cell.font = Font(italic=True, color="475569")

        ws.row_dimensions[row].height = min(120, max(30, len(a["content"]) // 5))

    # Freeze header
    ws.freeze_panes = "A2"

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def to_pptx(messages: list[dict]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(0x1B, 0x3A, 0x6B)
    GOLD = RGBColor(0xC8, 0xA8, 0x4B)
    DARK = RGBColor(0x1E, 0x29, 0x3B)
    GREY = RGBColor(0x64, 0x74, 0x8B)

    blank_layout = prs.slide_layouts[6]  # blank

    def add_rect(slide, l, t, w, h, color):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        return shape

    def add_textbox(slide, text, l, t, w, h, bold=False, size=18, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(size)
        run.font.color.rgb = color
        return txBox

    # ---- Title slide ----
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 5.8, 13.33, 0.25, GOLD)

    add_textbox(slide, "Governance Assistant", 1, 2.2, 11, 1.2,
                bold=True, size=40, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(slide, "Conversation Export", 1, 3.5, 11, 0.7,
                size=22, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)
    add_textbox(slide, datetime.now().strftime("%d %B %Y"), 1, 4.3, 11, 0.5,
                size=14, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)

    # ---- Q&A slides ----
    pairs = _qa_pairs(messages)
    MAX_ANSWER_CHARS = 900

    for i, (q, a) in enumerate(pairs, 1):
        slide = prs.slides.add_slide(blank_layout)

        # Top bar
        add_rect(slide, 0, 0, 13.33, 0.12, NAVY)
        # Bottom bar
        add_rect(slide, 0, 7.3, 13.33, 0.2, GOLD)

        # Q label pill
        add_rect(slide, 0.4, 0.25, 0.55, 0.42, NAVY)
        add_textbox(slide, f"Q{i}", 0.4, 0.25, 0.55, 0.42,
                    bold=True, size=13, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

        # Question
        add_textbox(slide, q["content"], 1.1, 0.22, 11.8, 0.9,
                    bold=True, size=16, color=NAVY)

        # Divider line
        add_rect(slide, 0.4, 1.25, 12.5, 0.03, RGBColor(0xE2, 0xE8, 0xF0))

        # Answer
        answer_text = a["content"]
        truncated = False
        if len(answer_text) > MAX_ANSWER_CHARS:
            answer_text = answer_text[:MAX_ANSWER_CHARS] + "…"
            truncated = True
        add_textbox(slide, answer_text, 0.4, 1.4, 12.5, 5.4, size=13, color=DARK)

        # Sources
        sources = a.get("sources", [])
        if sources:
            src_text = "Sources: " + "  ·  ".join(sources)
            add_textbox(slide, src_text, 0.4, 6.85, 10, 0.35,
                        size=9, color=GREY)

        # Slide number
        add_textbox(slide, f"{i} / {len(pairs)}", 12.5, 6.85, 0.8, 0.35,
                    size=9, color=GREY, align=PP_ALIGN.RIGHT)

        if truncated:
            add_textbox(slide, "* Full answer available in Word or PDF export", 0.4, 7.1, 8, 0.25,
                        size=8, color=GREY)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def to_pdf(messages: list[dict]) -> bytes:
    from fpdf import FPDF

    NAVY = (27, 58, 107)
    GOLD = (200, 168, 75)
    DARK = (30, 41, 59)
    GREY = (100, 116, 139)
    LIGHT = (239, 246, 255)

    class GovPDF(FPDF):
        def header(self):
            self.set_fill_color(*NAVY)
            self.rect(0, 0, 210, 8, "F")
            self.set_fill_color(*GOLD)
            self.rect(0, 8, 210, 1.5, "F")

        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(*GREY)
            self.cell(0, 10, f"Governance Assistant Export  ·  Page {self.page_no()}", align="C")

    pdf = GovPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()

    # Title block
    pdf.set_xy(15, 18)
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(*NAVY)
    pdf.cell(0, 10, "Governance Assistant", ln=True, align="C")

    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(*GREY)
    pdf.cell(0, 6, f"Conversation Export  ·  {datetime.now().strftime('%d %B %Y, %H:%M')}", ln=True, align="C")
    pdf.ln(8)

    pairs = _qa_pairs(messages)
    for i, (q, a) in enumerate(pairs, 1):
        # Q label + question
        pdf.set_fill_color(*LIGHT)
        pdf.set_draw_color(*NAVY)
        pdf.set_line_width(0.5)
        x, y = pdf.get_x(), pdf.get_y()

        # Q badge
        pdf.set_fill_color(*NAVY)
        pdf.set_font("Helvetica", "B", 9)
        pdf.set_text_color(255, 255, 255)
        pdf.rect(15, y, 10, 7, "F")
        pdf.set_xy(15, y + 0.5)
        pdf.cell(10, 6, f"Q{i}", align="C")

        # Question text
        pdf.set_font("Helvetica", "B", 11)
        pdf.set_text_color(*NAVY)
        pdf.set_xy(27, y)
        pdf.multi_cell(168, 7, q["content"])
        pdf.ln(2)

        # Answer
        pdf.set_font("Helvetica", "", 10)
        pdf.set_text_color(*DARK)
        pdf.set_x(18)
        pdf.multi_cell(177, 6, a["content"])
        pdf.ln(1)

        # Sources
        sources = a.get("sources", [])
        if sources:
            pdf.set_font("Helvetica", "I", 8)
            pdf.set_text_color(*GREY)
            pdf.set_x(18)
            pdf.multi_cell(177, 5, "Sources: " + "  ·  ".join(sources))

        # Separator
        pdf.ln(3)
        pdf.set_draw_color(226, 232, 240)
        pdf.set_line_width(0.3)
        pdf.line(15, pdf.get_y(), 195, pdf.get_y())
        pdf.ln(5)

    return bytes(pdf.output())
