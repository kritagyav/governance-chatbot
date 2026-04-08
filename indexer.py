"""
Document indexer for governance chatbot.
Uses TF-IDF search (no neural model) — lightweight, runs on 512MB RAM.
Handles PDF, DOCX, PPTX, XLSX, XLSM.
"""

import os
import hashlib
import json
import pickle
from pathlib import Path
from typing import Optional

import gc
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

import pypdf
from docx import Document as DocxDocument
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".ppt", ".xlsx", ".xls", ".xlsm"}
CHUNK_WORDS = 400
CHUNK_OVERLAP_WORDS = 80


class GovernanceIndexer:
    def __init__(self, folder_path: str, db_path: str = "./search_db"):
        self.folder_path = str(folder_path)
        os.makedirs(db_path, exist_ok=True)

        self._chunks_path  = os.path.join(db_path, "chunks.json")
        self._hashes_path  = os.path.join(db_path, "hashes.json")
        self._tfidf_path   = os.path.join(db_path, "tfidf.pkl")

        # chunks: list of {"text": str, "source": str, "path": str}
        self._chunks  = self._load_json(self._chunks_path, [])
        self._hashes  = self._load_json(self._hashes_path, {})
        self._vectorizer: Optional[TfidfVectorizer] = None
        self._matrix   = None
        self._load_tfidf()

    # ------------------------------------------------------------------ #
    # Persistence helpers
    # ------------------------------------------------------------------ #

    def _load_json(self, path, default):
        if os.path.exists(path):
            with open(path) as f:
                return json.load(f)
        return default

    def _save_json(self, path, obj):
        with open(path, "w") as f:
            json.dump(obj, f)

    def _load_tfidf(self):
        if os.path.exists(self._tfidf_path) and self._chunks:
            try:
                with open(self._tfidf_path, "rb") as f:
                    self._vectorizer, self._matrix = pickle.load(f)
            except Exception:
                self._rebuild_tfidf()

    def _rebuild_tfidf(self):
        if not self._chunks:
            self._vectorizer = None
            self._matrix = None
            return
        texts = [c["text"] for c in self._chunks]
        self._vectorizer = TfidfVectorizer(
            max_features=20000,
            ngram_range=(1, 2),
            sublinear_tf=True,
        )
        self._matrix = self._vectorizer.fit_transform(texts)
        with open(self._tfidf_path, "wb") as f:
            pickle.dump((self._vectorizer, self._matrix), f)

    # ------------------------------------------------------------------ #
    # Document loaders
    # ------------------------------------------------------------------ #

    def _load_pdf(self, path: str) -> str:
        reader = pypdf.PdfReader(path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    def _load_docx(self, path: str) -> str:
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _load_pptx(self, path: str) -> str:
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    def _load_xlsx(self, path: str) -> str:
        """Load xlsx/xlsm using openpyxl in read_only+data_only mode (memory efficient)."""
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True, max_row=500):
                row_data = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in row_data):
                    rows.append("  |  ".join(row_data))
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)

    def _load_xls(self, path: str) -> str:
        """Load legacy .xls using xlrd directly (no pandas needed)."""
        import xlrd
        wb = xlrd.open_workbook(path)
        parts = []
        for sheet in wb.sheets():
            rows = []
            for row_idx in range(min(sheet.nrows, 500)):
                row_data = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                if any(c.strip() for c in row_data):
                    rows.append("  |  ".join(row_data))
            if rows:
                parts.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
        return "\n\n".join(parts)

    def _load_document(self, path: str) -> Optional[str]:
        ext = Path(path).suffix.lower()
        try:
            if ext == ".pdf":
                return self._load_pdf(path)
            elif ext == ".docx":
                return self._load_docx(path)
            elif ext in (".pptx", ".ppt"):
                return self._load_pptx(path)
            elif ext == ".xls":
                return self._load_xls(path)
            elif ext in (".xlsx", ".xlsm"):
                return self._load_xlsx(path)
        except Exception as e:
            print(f"[Indexer] Error loading {Path(path).name}: {e}")
        return None

    # ------------------------------------------------------------------ #
    # Chunking
    # ------------------------------------------------------------------ #

    def _chunk_text(self, text: str) -> list[str]:
        words = text.split()
        step = CHUNK_WORDS - CHUNK_OVERLAP_WORDS
        chunks = []
        for i in range(0, len(words), step):
            chunk = " ".join(words[i: i + CHUNK_WORDS])
            if len(chunk.strip()) > 50:
                chunks.append(chunk)
        return chunks

    # ------------------------------------------------------------------ #
    # Index operations
    # ------------------------------------------------------------------ #

    def _file_hash(self, path: str) -> str:
        with open(path, "rb") as f:
            return hashlib.md5(f.read()).hexdigest()

    def _remove_file_chunks(self, path: str):
        before = len(self._chunks)
        self._chunks = [c for c in self._chunks if c["path"] != path]
        if len(self._chunks) != before:
            self._save_json(self._chunks_path, self._chunks)

    def index_file(self, path: str) -> bool:
        if Path(path).suffix.lower() not in SUPPORTED_EXTENSIONS:
            return False
        if not os.path.exists(path):
            return False

        current_hash = self._file_hash(path)
        if self._hashes.get(path) == current_hash:
            return False  # unchanged

        self._remove_file_chunks(path)

        text = self._load_document(path)
        if not text or not text.strip():
            print(f"[Indexer] No text extracted from {Path(path).name}")
            return False

        chunks = self._chunk_text(text)
        if not chunks:
            return False

        filename = Path(path).name
        for chunk in chunks:
            self._chunks.append({"text": chunk, "source": filename, "path": path})

        self._hashes[path] = current_hash
        self._save_json(self._chunks_path, self._chunks)
        self._save_json(self._hashes_path, self._hashes)
        self._rebuild_tfidf()
        gc.collect()
        print(f"[Indexer] Indexed '{filename}' → {len(chunks)} chunks")
        return True

    def remove_file(self, path: str):
        self._remove_file_chunks(path)
        self._hashes.pop(path, None)
        self._save_json(self._hashes_path, self._hashes)
        self._rebuild_tfidf()
        print(f"[Indexer] Removed '{Path(path).name}'")

    def index_all(self):
        folder = Path(self.folder_path)
        if not folder.exists():
            return

        files = [f for f in folder.rglob("*")
                 if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS]
        print(f"[Indexer] Scanning {len(files)} document(s)…")

        updated = sum(self.index_file(str(f)) for f in files)

        # Remove stale entries
        current_paths = {str(f) for f in files}
        stale = [p for p in list(self._hashes) if p not in current_paths]
        for p in stale:
            self._remove_file_chunks(p)
            del self._hashes[p]
        if stale:
            self._save_json(self._hashes_path, self._hashes)
            self._rebuild_tfidf()

        print(f"[Indexer] Done — {updated} new/updated, {len(stale)} removed.")

    def force_reindex_all(self):
        self._hashes.clear()
        self._save_json(self._hashes_path, self._hashes)
        self.index_all()

    # ------------------------------------------------------------------ #
    # Search
    # ------------------------------------------------------------------ #

    def search(self, query: str, n_results: int = 6) -> list[dict]:
        if not self._chunks or self._vectorizer is None:
            return []

        query_vec = self._vectorizer.transform([query])
        scores = cosine_similarity(query_vec, self._matrix)[0]
        top_indices = np.argsort(scores)[::-1][:n_results]

        return [
            {
                "text": self._chunks[i]["text"],
                "source": self._chunks[i]["source"],
                "path": self._chunks[i]["path"],
                "score": float(scores[i]),
            }
            for i in top_indices
            if scores[i] > 0
        ]

    # ------------------------------------------------------------------ #
    # Status
    # ------------------------------------------------------------------ #

    def get_indexed_files(self) -> list[str]:
        return sorted({Path(p).name for p in self._hashes})

    def get_chunk_count(self) -> int:
        return len(self._chunks)
